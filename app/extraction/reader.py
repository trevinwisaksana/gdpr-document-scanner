from __future__ import annotations

import csv
import io
from pathlib import Path


def extract_text(file_path: str | Path) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _read_pdf(path)
    elif suffix == ".docx":
        return _read_docx(path)
    elif suffix == ".pptx":
        return _read_pptx(path)
    elif suffix in (".xls", ".xlsx"):
        return _read_excel(path)
    elif suffix == ".csv":
        return _read_csv(path)
    elif suffix in (".html", ".htm"):
        return _read_html(path)
    elif suffix == ".rtf":
        return _read_rtf(path)
    elif suffix in (".txt", ".md", ".log", ".json", ".xml", ".yaml", ".yml"):
        return path.read_text(encoding="utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: {suffix!r}")


def _read_pdf(path: Path) -> str:
    import pymupdf

    doc = pymupdf.open(path)
    return "\n".join(page.get_text() for page in doc)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
    return "\n".join(texts)


def _read_excel(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            lines.append("\t".join("" if v is None else str(v) for v in row))
    return "\n".join(lines)


def _read_csv(path: Path) -> str:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        return "\n".join("\t".join(row) for row in reader)


def _read_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    html = path.read_text(encoding="utf-8", errors="replace")
    return BeautifulSoup(html, "html.parser").get_text(separator="\n")


def _read_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    raw = path.read_text(encoding="utf-8", errors="replace")
    return rtf_to_text(raw)
