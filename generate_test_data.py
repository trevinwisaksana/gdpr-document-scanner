"""
Generate synthetic PII-laden test documents and upload to Google Drive.

Usage:
    python generate_test_data.py --folder-id <DRIVE_FOLDER_ID>

Auth (choose one):
    --oauth   Use OAuth user credentials (opens browser, works with personal Gmail).
              Requires client_secrets.json (OAuth Desktop app from GCP Console).
    default   Uses service account credentials.json (requires a Shared Drive).
"""

import argparse
import csv
import io
import json
import os
import random
import tempfile

import fitz  # PyMuPDF
import openpyxl
from docx import Document
from docx.shared import Pt
from google.oauth2 import service_account
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

CREDENTIALS_FILE = os.path.join(os.path.dirname(__file__), "credentials.json")
CLIENT_SECRETS_FILE = os.path.join(os.path.dirname(__file__), "client_secrets.json")
TOKEN_FILE = os.path.join(os.path.dirname(__file__), "token.json")
SCOPES = ["https://www.googleapis.com/auth/drive"]

# --- Synthetic PII data ---

RECORDS = [
    {
        "name": "Alice Johnson",
        "email": "alice.johnson@example.com",
        "phone": "+1-415-555-0192",
        "ssn": "123-45-6789",
        "dob": "1985-03-14",
        "address": "742 Evergreen Terrace, Springfield, IL 62701",
        "credit_card": "4111 1111 1111 1111",
        "passport": "US-A12345678",
        "ip": "192.168.1.42",
        "medical_id": "MRN-00482910",
    },
    {
        "name": "Bob Martinez",
        "email": "b.martinez@healthco.org",
        "phone": "(312) 555-7834",
        "ssn": "987-65-4321",
        "dob": "1972-11-30",
        "address": "19 Maple Ave, Chicago, IL 60601",
        "credit_card": "5500 0000 0000 0004",
        "passport": "GB-B98765432",
        "ip": "10.0.0.7",
        "medical_id": "MRN-00193847",
    },
    {
        "name": "Clara Schmidt",
        "email": "clara.schmidt@fintech.de",
        "phone": "+49 30 12345678",
        "ssn": "456-78-9012",
        "dob": "1990-07-22",
        "address": "Unter den Linden 5, 10117 Berlin, Germany",
        "credit_card": "3714 496353 98431",
        "passport": "DE-C11223344",
        "ip": "172.16.254.1",
        "medical_id": "MRN-00567234",
    },
    {
        "name": "David Chen",
        "email": "dchen@acmecorp.io",
        "phone": "+1-650-555-0181",
        "ssn": "321-54-9876",
        "dob": "1968-01-05",
        "address": "1600 Amphitheatre Pkwy, Mountain View, CA 94043",
        "credit_card": "6011 1111 1111 1117",
        "passport": "US-D44556677",
        "ip": "203.0.113.45",
        "medical_id": "MRN-00871029",
    },
    {
        "name": "Elena Rossi",
        "email": "elena.rossi@studio.it",
        "phone": "+39 06 1234 5678",
        "ssn": "654-32-1098",
        "dob": "1995-09-18",
        "address": "Via Roma 12, 00185 Rome, Italy",
        "credit_card": "4012 8888 8888 1881",
        "passport": "IT-E55667788",
        "ip": "198.51.100.23",
        "medical_id": "MRN-00334456",
    },
]

REPORT_TEMPLATE = """\
CONFIDENTIAL PERSONNEL RECORD
==============================
Name:           {name}
Date of Birth:  {dob}
SSN:            {ssn}
Passport:       {passport}

Contact Information
-------------------
Email:          {email}
Phone:          {phone}
Address:        {address}
IP Address:     {ip}

Financial & Medical
--------------------
Credit Card:    {credit_card}
Medical ID:     {medical_id}

This document contains personally identifiable information (PII) and is
subject to GDPR data protection regulations. Unauthorized disclosure is
strictly prohibited.
"""


# --- Document generators ---

def make_txt(record: dict) -> bytes:
    return REPORT_TEMPLATE.format(**record).encode()


def make_pdf(record: dict) -> bytes:
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)  # A4
    text = REPORT_TEMPLATE.format(**record)
    page.insert_text((60, 60), text, fontsize=11, fontname="helv")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_docx(record: dict) -> bytes:
    doc = Document()
    doc.add_heading("Confidential Personnel Record", level=1)
    for key, value in record.items():
        p = doc.add_paragraph()
        run = p.add_run(f"{key.replace('_', ' ').title()}: ")
        run.bold = True
        run.font.size = Pt(11)
        p.add_run(value).font.size = Pt(11)
    doc.add_paragraph(
        "\nThis document contains PII subject to GDPR data protection regulations."
    )
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_csv(records: list[dict]) -> bytes:
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=records[0].keys())
    writer.writeheader()
    writer.writerows(records)
    return buf.getvalue().encode()


def make_xlsx(records: list[dict]) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Employee PII"
    headers = list(records[0].keys())
    ws.append([h.replace("_", " ").title() for h in headers])
    for rec in records:
        ws.append([rec[h] for h in headers])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx(records: list[dict]) -> bytes:
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # title + content

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Employee PII Data — Confidential"
    title_slide.placeholders[1].text = "GDPR Scanner Test Dataset"

    for rec in records:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = rec["name"]
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        fields = ["email", "phone", "ssn", "dob", "address", "credit_card", "passport", "ip", "medical_id"]
        for field in fields:
            p = tf.add_paragraph()
            p.text = f"{field.replace('_', ' ').title()}: {rec[field]}"
            p.font.size = PptPt(16)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Realistic subfolder names to simulate a real Drive structure
SUBFOLDERS = [
    "HR / Onboarding",
    "HR / Payroll Records",
    "Finance / Invoices 2024",
    "Finance / Expense Reports",
    "Legal / Contracts",
    "Legal / Compliance",
    "IT / Access Requests",
    "Operations / Incident Reports",
    "Medical / Patient Files",
    "Marketing / Campaign Data",
]


# --- Drive upload ---

def get_drive_service():
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request

    token_path = os.path.join(os.path.dirname(__file__), "drive_token.json")

    creds = None
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CLIENT_SECRETS_FILE, SCOPES)
            auth_url, _ = flow.authorization_url(prompt="consent")
            print(f"\nOpen this URL in your browser:\n\n{auth_url}\n")
            import subprocess
            subprocess.run(["open", auth_url])
            creds = flow.run_local_server(port=8080, open_browser=False)
        with open(token_path, "w") as f:
            f.write(creds.to_json())

    return build("drive", "v3", credentials=creds)


def create_folder(service, name: str, parent_id: str) -> str:
    metadata = {
        "name": name,
        "mimeType": "application/vnd.google-apps.folder",
        "parents": [parent_id],
    }
    folder = service.files().create(body=metadata, fields="id,name").execute()
    print(f"  Created folder: {folder['name']} (id={folder['id']})")
    return folder["id"]


def upload_file(service, folder_id: str, name: str, content: bytes, mime_type: str):
    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(name)[1]) as f:
        f.write(content)
        tmp_path = f.name

    try:
        metadata = {"name": name, "parents": [folder_id]}
        media = MediaFileUpload(tmp_path, mimetype=mime_type, resumable=False)
        file = (
            service.files()
            .create(body=metadata, media_body=media, fields="id,name")
            .execute()
        )
        print(f"    Uploaded: {file['name']} (id={file['id']})")
    finally:
        os.unlink(tmp_path)


def main():
    parser = argparse.ArgumentParser(description="Generate and upload PII test documents to Google Drive")
    parser.add_argument("--folder-id", required=True, help="Google Drive folder ID to upload into")
    parser.add_argument("--seed", type=int, default=42, help="Random seed for folder assignment")
    args = parser.parse_args()

    random.seed(args.seed)

    print("Connecting to Google Drive...")
    service = get_drive_service()

    files = [
        *[
            (f"employee_{i+1}_{rec['name'].split()[0].lower()}.txt", make_txt(rec), "text/plain")
            for i, rec in enumerate(RECORDS)
        ],
        *[
            (f"employee_{i+1}_{rec['name'].split()[0].lower()}.pdf", make_pdf(rec), "application/pdf")
            for i, rec in enumerate(RECORDS)
        ],
        *[
            (f"employee_{i+1}_{rec['name'].split()[0].lower()}.docx", make_docx(rec), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            for i, rec in enumerate(RECORDS)
        ],
        ("all_employees.csv", make_csv(RECORDS), "text/csv"),
        ("all_employees.xlsx", make_xlsx(RECORDS), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        ("employee_presentation.pptx", make_pptx(RECORDS), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ]

    # Randomly pick which subfolders to actually create (not necessarily all of them)
    chosen_names = random.sample(SUBFOLDERS, k=min(len(SUBFOLDERS), max(3, len(files) // 3)))

    print(f"\nCreating {len(chosen_names)} subfolders...\n")
    folder_ids = {name: create_folder(service, name, args.folder_id) for name in chosen_names}
    folder_list = list(folder_ids.values())

    print(f"\nUploading {len(files)} files into random subfolders...\n")
    for name, content, mime in files:
        target_folder_id = random.choice(folder_list)
        target_name = next(k for k, v in folder_ids.items() if v == target_folder_id)
        print(f"  -> {target_name}/")
        upload_file(service, target_folder_id, name, content, mime)

    print(f"\nDone. {len(files)} files spread across {len(folder_ids)} folders.")


if __name__ == "__main__":
    main()
