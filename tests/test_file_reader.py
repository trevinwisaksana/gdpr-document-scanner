import csv
import io
import pytest
from pathlib import Path

from app.extraction.reader import extract_text

SAMPLE = "John Doe lives at 123 Main St and his email is john@example.com"


# ── fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture
def pdf_file(tmp_path):
    import pymupdf

    path = tmp_path / "sample.pdf"
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), SAMPLE)
    doc.save(path)
    doc.close()
    return path


@pytest.fixture
def docx_file(tmp_path):
    from docx import Document

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph(SAMPLE)
    doc.save(path)
    return path


@pytest.fixture
def pptx_file(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    txBox.text_frame.text = SAMPLE
    prs.save(path)
    return path


@pytest.fixture
def xlsx_file(tmp_path):
    import openpyxl

    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Address", "Email"])
    ws.append(["John Doe", "123 Main St", "john@example.com"])
    wb.save(path)
    return path


@pytest.fixture
def csv_file(tmp_path):
    path = tmp_path / "sample.csv"
    with open(path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Address", "Email"])
        writer.writerow(["John Doe", "123 Main St", "john@example.com"])
    return path


@pytest.fixture
def html_file(tmp_path):
    path = tmp_path / "sample.html"
    path.write_text(f"<html><body><p>{SAMPLE}</p></body></html>", encoding="utf-8")
    return path


@pytest.fixture
def rtf_file(tmp_path):
    path = tmp_path / "sample.rtf"
    # Minimal valid RTF wrapping plain text
    path.write_text(r"{\rtf1\ansi " + SAMPLE + r"}", encoding="utf-8")
    return path


@pytest.fixture
def txt_file(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text(SAMPLE, encoding="utf-8")
    return path


# ── tests ─────────────────────────────────────────────────────────────────────

def test_extract_pdf(pdf_file):
    text = extract_text(pdf_file)
    assert "John Doe" in text
    assert "john@example.com" in text


def test_extract_docx(docx_file):
    text = extract_text(docx_file)
    assert SAMPLE in text


def test_extract_pptx(pptx_file):
    text = extract_text(pptx_file)
    assert "John Doe" in text
    assert "john@example.com" in text


def test_extract_xlsx(xlsx_file):
    text = extract_text(xlsx_file)
    assert "John Doe" in text
    assert "john@example.com" in text


def test_extract_csv(csv_file):
    text = extract_text(csv_file)
    assert "John Doe" in text
    assert "john@example.com" in text


def test_extract_html(html_file):
    text = extract_text(html_file)
    assert SAMPLE in text


def test_extract_rtf(rtf_file):
    text = extract_text(rtf_file)
    assert "John Doe" in text


def test_extract_txt(txt_file):
    text = extract_text(txt_file)
    assert SAMPLE in text


def test_extract_returns_string(txt_file):
    assert isinstance(extract_text(txt_file), str)


def test_unsupported_extension_raises(tmp_path):
    path = tmp_path / "file.xyz"
    path.write_text("data")
    with pytest.raises(ValueError, match="Unsupported file type"):
        extract_text(path)


def test_accepts_string_path(txt_file):
    text = extract_text(str(txt_file))
    assert SAMPLE in text
